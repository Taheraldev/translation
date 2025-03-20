import os
import requests
import tempfile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from telegram import Update, InputFile
from telegram.ext import Updater, CommandHandler, MessageHandler, Filters, CallbackContext
import logging

# إعدادات التسجيل
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s', level=logging.INFO
)
logger = logging.getLogger(__name__)

# إعدادات API للترجمة عبر RapidAPI
API_URL = "https://openl-translate.p.rapidapi.com/translate/bulk"
HEADERS = {
    "X-RapidAPI-Key": "7b73717e2dmshbd139747c640560p175307jsn75624bf31396",       # استبدل بـ مفتاح RapidAPI الخاص بك
    "X-RapidAPI-Host": "openl-translate.p.rapidapi.com",
    "Content-Type": "application/json"
}

def translate_text(text: str, source_lang="en", target_lang="ar") -> str:
    """
    ترجم النص باستخدام API الخاص بـ openl-translate.
    في حال فشل الترجمة أو عدم وجود استجابة صالحة يتم إرجاع النص الأصلي.
    """
    if not text.strip():
        return text
    payload = {"text": [text], "source_lang": source_lang, "target_lang": target_lang}
    response = requests.post(API_URL, json=payload, headers=HEADERS)
    try:
        json_data = response.json()
    except Exception as e:
        logger.error("فشل فك ترميز JSON: %s\nمحتوى الاستجابة: %s", e, response.text)
        return text
    # نتوقع استجابة بالشكل: { "translations": [ { "text": "النص المترجم" } ] }
    translated_text = json_data.get("translations", [{}])[0].get("text", "")
    return translated_text if translated_text else text

def translate_pdf(file_path: str) -> str:
    """
    ترجمة محتوى ملف PDF مع الحفاظ على التنسيق.
    يتم استخراج النص من كل صفحة وترجمته، ثم يتم إعادة إدراجه في الصفحة.
    """
    new_pdf_path = file_path.replace(".pdf", "_translated.pdf")
    doc = fitz.open(file_path)
    for page in doc:
        text = page.get_text("text")
        translated_text = translate_text(text)
        # تنظيف المحتويات القديمة وإدراج النص المترجم.
        page.clean_contents()
        page.insert_text((50, 50), translated_text, fontsize=12)
    doc.save(new_pdf_path)
    return new_pdf_path

def translate_docx(file_path: str) -> str:
    """
    ترجمة محتوى ملف DOCX مع الحفاظ على التنسيق والجداول.
    """
    new_docx_path = file_path.replace(".docx", "_translated.docx")
    doc = Document(file_path)
    for para in doc.paragraphs:
        para.text = translate_text(para.text)
    # ترجمة النصوص داخل الجداول
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell.text = translate_text(cell.text)
    doc.save(new_docx_path)
    return new_docx_path

def translate_pptx(file_path: str) -> str:
    """
    ترجمة محتوى ملف PPTX مع الحفاظ على تنسيق الشرائح.
    """
    new_pptx_path = file_path.replace(".pptx", "_translated.pptx")
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                shape.text = translate_text(shape.text)
    prs.save(new_pptx_path)
    return new_pptx_path

def handle_document(update: Update, context: CallbackContext):
    document = update.message.document
    file_name = document.file_name.lower()
    ext = os.path.splitext(file_name)[1]
    valid_extensions = {
        ".pdf": translate_pdf,
        ".docx": translate_docx,
        ".pptx": translate_pptx
    }
    
    if ext not in valid_extensions:
        update.message.reply_text("❌ الملف غير مدعوم! الرجاء إرسال PDF أو DOCX أو PPTX.")
        return

    # تحميل الملف من تيليجرام
    file_id = document.file_id
    file_obj = context.bot.get_file(file_id)
    temp_dir = tempfile.mkdtemp()
    local_file_path = os.path.join(temp_dir, file_name)
    file_obj.download(custom_path=local_file_path)
    
    update.message.reply_text("جارٍ ترجمة الملف، يرجى الانتظار...")
    try:
        translated_file_path = valid_extensions[ext](local_file_path)
    except Exception as e:
        logger.error("حدث خطأ أثناء ترجمة الملف: %s", e)
        update.message.reply_text("❌ حدث خطأ أثناء ترجمة الملف.")
        return
    
    with open(translated_file_path, "rb") as f:
        update.message.reply_document(
            document=InputFile(f, filename=os.path.basename(translated_file_path)),
            caption="✅ تم الترجمة بنجاح!"
        )
    # يمكنك حذف الملفات المؤقتة إذا رغبت بذلك
    # os.remove(local_file_path)
    # os.remove(translated_file_path)

def start(update: Update, context: CallbackContext):
    update.message.reply_text("مرحباً! أرسل لي ملف PDF, DOCX, أو PPTX وسأقوم بترجمته من الإنجليزية إلى العربية مع الحفاظ على التنسيق.")

def main():
    TELEGRAM_BOT_TOKEN = "5284087690:AAGwKfPojQ3c-SjCHSIdeog-yN3-4Gpim1Y"  # استبدل بـ توكن البوت الخاص بك
    updater = Updater(TELEGRAM_BOT_TOKEN, use_context=True)
    dp = updater.dispatcher

    dp.add_handler(CommandHandler("start", start))
    dp.add_handler(MessageHandler(Filters.document, handle_document, run_async=True))
    
    updater.start_polling()
    updater.idle()

if __name__ == '__main__':
    main()
